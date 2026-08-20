r"""
로컬에서 실행하는 문서 등록 프로그램 (이미지 추출 + OCR 지원). Windows/macOS/Linux 공용 코드.
- 인자 없이 실행 -> 이 스크립트가 있는 폴더의 "incoming" 하위 폴더를 통째로 등록
- 특정 파일/폴더를 지정하면 그것만 등록
- PDF 안의 이미지는 별도 파일로 추출 저장 + OCR 텍스트를 함께 Qdrant에 등록

지원 파일 형식:
    PDF(.pdf), Word(.docx/.doc), 한글(.hwp/.hwpx), 엑셀(.xlsx),
    파워포인트(.pptx/.ppt), 텍스트(.txt/.md), HTML(.html/.htm), 이미지(.jpg/.jpeg/.png)

사용법:
    python register.py
    python register.py "문서.pdf"
    python register.py "C:\Docs\등록할문서들"      (Windows)
    python register.py "/Users/me/Docs/등록할문서들"  (macOS/Linux)

설정:
    이 스크립트와 같은 폴더의 config.json 에서 mcp_url 등을 읽어옵니다.
    파일이 없으면 최초 실행 시 기본값으로 자동 생성됩니다.

사전 설치:
    pip install pymupdf python-docx pyhwp mcp pytesseract pillow requests openpyxl python-pptx

OCR을 쓰려면 Tesseract 엔진 자체도 설치해야 합니다(파이썬 패키지와는 별도).
PATH에 있으면 자동으로 찾고, 없으면 config.json의 tesseract_cmd에 직접 경로를 지정하세요:
    - Windows: https://github.com/UB-Mannheim/tesseract/wiki 에서 설치파일 다운로드
      (설치 중 "Additional language data"에서 Korean 체크, 보통
      C:\Program Files\Tesseract-OCR\tesseract.exe 에 설치됨)
    - macOS: brew install tesseract tesseract-lang
    - Linux(Debian/Ubuntu 계열): sudo apt install tesseract-ocr tesseract-ocr-kor
    Tesseract가 없어도 스크립트는 동작합니다 - 이 경우 이미지는 저장되지만
    OCR 텍스트 없이 "페이지/파일명"만으로 등록됩니다.

이미지 "의미" 캡션(비전 모델)을 쓰려면:
    1) Ollama 설치: https://ollama.com/download
       (Windows: winget install Ollama.Ollama / macOS: brew install ollama)
    2) 모델 다운로드: ollama pull moondream
       (CPU 전용 PC 기준 가벼운 모델. GPU 있으면 llava, qwen2.5vl 등으로 교체 가능)
    Ollama가 꺼져 있거나 모델이 없어도 스크립트는 동작합니다 - 이 경우 캡션 없이
    OCR 텍스트(있으면)만으로 등록됩니다.

구버전 .doc/.ppt 지원(선택):
    antiword(.doc) 또는 LibreOffice(둘 다 변환용)가 설치되어 있으면 자동 사용됩니다.
    - macOS: brew install --cask libreoffice
    - Linux: sudo apt install libreoffice
    없으면 해당 파일은 건너뛰고 안내 메시지를 출력합니다.
"""
import asyncio
import base64
import io
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from pathlib import Path

import certifi

# mcp(httpx2)가 기본으로 쓰는 truststore가 Python 3.14 + Windows 조합에서
# SSLContext.verify_mode 설정 시 무한 재귀에 빠지는 버그가 있음. SSL_CERT_FILE을
# 지정해 truststore 대신 certifi 인증서로 검증하도록 우회. (mcp import 전에 설정해야 함)
os.environ.setdefault("SSL_CERT_FILE", certifi.where())

import docx
import fitz  # PyMuPDF
import openpyxl
import requests
from mcp import ClientSession
from mcp.client.streamable_http import streamable_http_client as streamablehttp_client
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class _SuppressSSETeardownNoise(logging.Filter):
    """mcp 라이브러리 자체 주석에 명시된 알려진 동작: 등록이 끝나고 연결을
    정리(취소)하는 순간 SSE 스트림에 마지막 이벤트가 도착하면 ClosedResourceError가
    나면서 'Error parsing SSE message'로 로그된다. 실제 데이터 등록에는 영향 없는
    타이밍성 노이즈라 이 메시지만 걸러낸다."""

    def filter(self, record: logging.LogRecord) -> bool:
        return record.getMessage() != "Error parsing SSE message"


logging.getLogger("mcp.client.streamable_http").addFilter(_SuppressSSETeardownNoise())

try:
    import pytesseract
    OCR_LIB_AVAILABLE = True
except ImportError:
    OCR_LIB_AVAILABLE = False

# ---------------------------------------------------------------
# 환경설정 (config.json에서 로드, 없으면 기본값으로 자동 생성)
# ---------------------------------------------------------------
if getattr(sys, "frozen", False):
    # PyInstaller onefile로 실행 중이면 __file__은 임시 압축 해제 폴더를 가리키므로,
    # config.json/incoming/extracted_images는 실제 exe가 있는 폴더 기준으로 잡는다.
    _exe_path = Path(sys.executable).resolve()
    if _exe_path.parent.name == "MacOS" and _exe_path.parent.parent.name == "Contents":
        # macOS .app 번들 안에서는 exe가 Contents/MacOS 깊숙이 있으므로,
        # 번들 내부가 아니라 .app이 놓인(사용자 눈에 보이는) 폴더를 기준으로 잡는다.
        # Windows exe와 동일하게 "실행 파일과 같은 경로"로 보이게 하기 위함.
        SCRIPT_DIR = _exe_path.parents[3]
    else:
        SCRIPT_DIR = _exe_path.parent
else:
    SCRIPT_DIR = Path(__file__).resolve().parent
CONFIG_PATH = SCRIPT_DIR / "config.json"

DEFAULT_CONFIG = {
    # 실제 서버 주소/키는 이 소스에 넣지 말고 config.json에서만 관리 (config.json은 git 추적 제외)
    "mcp_url": "https://example.com/mcp?key=REPLACE_ME",
    "ollama_url": "http://localhost:11434/api/generate",
    "vision_model": "moondream",  # GPU 있으면 llava, qwen2.5vl 등으로 교체 가능
    "tesseract_cmd": r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    "process_images": "y",  # y: 이미지 추출/OCR/캡션 처리, n: 이미지는 건너뛰고 텍스트만 등록
    "store_image_base64": "y",  # y: 이미지 파일 자체를 base64로 인코딩해 Qdrant metadata에 함께 저장
}

MAX_BASE64_IMAGE_BYTES = 5 * 1024 * 1024  # 이보다 큰 이미지는 용량 문제로 base64 저장 생략


def parse_yn(value, default: bool = True) -> bool:
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in ("y", "yes", "true", "1")


def load_config() -> dict:
    if CONFIG_PATH.exists():
        try:
            data = json.loads(CONFIG_PATH.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"config.json 읽기 실패, 기본값 사용: {e}")
            data = {}
    else:
        data = {}
        CONFIG_PATH.write_text(json.dumps(DEFAULT_CONFIG, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"config.json이 없어 기본값으로 생성했습니다: {CONFIG_PATH}")
    return {**DEFAULT_CONFIG, **data}


CONFIG = load_config()
MCP_URL = CONFIG["mcp_url"]
OLLAMA_URL = CONFIG["ollama_url"]
VISION_MODEL = CONFIG["vision_model"]
PROCESS_IMAGES = parse_yn(CONFIG["process_images"])
STORE_IMAGE_BASE64 = parse_yn(CONFIG["store_image_base64"])

CHUNK_SIZE = 1500  # 800 -> 1500: 청크 수 자체를 줄여 저장 요청/배치 수를 더 줄임 (검색 정밀도 손해는 미미)
CHUNK_OVERLAP = 200

CAPTION_PROMPT = "Describe what this image shows and what information or meaning it conveys, in detail."
# 참고: moondream은 한국어 출력 품질이 낮아(예: 한 단어만 반환) 캡션은 영어로 생성합니다.
# 이미지 속 한국어 원문은 OCR(pytesseract, kor+eng)이 별도로 처리합니다.

DEFAULT_INCOMING_DIR = SCRIPT_DIR / "incoming"
IMAGES_DIR = SCRIPT_DIR / "extracted_images"
EXTRACTED_TEXT_DIR = SCRIPT_DIR / "extracted_text"  # Qdrant에 실제로 보낸 내용을 파일별로 남기는 폴더
IMAGE_EXTS = {".jpg", ".jpeg", ".png"}


def find_soffice() -> str | None:
    """구버전 .doc/.ppt 변환에 쓸 LibreOffice(soffice) 실행 파일을 찾는다 (Windows/macOS/Linux)."""
    exe = shutil.which("soffice") or shutil.which("soffice.bin")
    if exe:
        return exe
    for candidate in (
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        # Linux (배포판/설치 방식에 따라 경로가 다양함)
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/opt/libreoffice/program/soffice",
        "/snap/bin/libreoffice.soffice",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def find_tesseract(configured: str) -> str | None:
    """Tesseract 실행 파일을 찾는다 (Windows/macOS/Linux). config.json에 지정된 경로가
    실제로 존재하면 그걸 쓰고, 아니면 PATH -> OS별 흔한 설치 경로 순으로 자동 탐색한다."""
    if configured and Path(configured).exists():
        return configured
    exe = shutil.which("tesseract")
    if exe:
        return exe
    for candidate in (
        # Windows
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        # macOS (Homebrew: Apple Silicon vs Intel)
        "/opt/homebrew/bin/tesseract",
        "/usr/local/bin/tesseract",
        # Linux
        "/usr/bin/tesseract",
    ):
        if Path(candidate).exists():
            return candidate
    return None


SOFFICE_CMD = find_soffice()

TESSERACT_CMD = find_tesseract(CONFIG["tesseract_cmd"])
if OCR_LIB_AVAILABLE and TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    OCR_AVAILABLE = True
else:
    OCR_AVAILABLE = False


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """줄(line) 경계를 지키면서 청크를 만든다. 예전에는 문자 수로만 뚝 잘라서(text[i:i+size])
    엑셀 행("셀1\\t셀2\\t셀3")이나 문장이 청크 중간에서 끊겨 검색 결과가 잘린 것처럼 보이는
    문제가 있었다 - 한 줄이 통째로 다음 청크로 넘어가더라도 절대 줄 중간을 자르지 않는다.
    단, 한 줄 자체가 size보다 긴 경우(줄바꿈 없는 매우 긴 문단 등)는 예외적으로 그 줄만
    문자 단위로 쪼갠다."""
    lines = []
    for line in text.split("\n"):
        if len(line) <= size:
            lines.append(line)
        else:
            for i in range(0, len(line), size):
                lines.append(line[i:i + size])

    chunks, i, n = [], 0, len(lines)
    while i < n:
        buf, length, j = [], 0, i
        while j < n:
            add_len = len(lines[j]) + (1 if buf else 0)
            if buf and length + add_len > size:
                break
            buf.append(lines[j])
            length += add_len
            j += 1
        chunks.append("\n".join(buf))
        if j >= n:
            break
        # 다음 청크는 문자 기준 overlap만큼 줄 단위로 되돌아가서 시작 (최소 한 줄은 전진)
        back_chars, k = 0, j
        while k > i and back_chars < overlap:
            k -= 1
            back_chars += len(lines[k]) + 1
        i = max(i + 1, k)
    return [c for c in chunks if c.strip()]


BASE64_MAX_SIDE = 1600  # Qdrant에 저장할 base64 이미지는 원본 대신 이 크기로 축소해서 용량 절약


def encode_image_base64(img_path: Path) -> str:
    """이미지를 (필요하면) 축소한 뒤 base64 문자열로 인코딩. 로컬에 저장된 원본 파일 자체는 그대로 둔다."""
    try:
        img = Image.open(img_path)
        fmt = img.format or "PNG"
        if max(img.size) > BASE64_MAX_SIDE:
            scale = BASE64_MAX_SIDE / max(img.size)
            new_size = (round(img.width * scale), round(img.height * scale))
            img = img.resize(new_size, Image.LANCZOS)

        buf = io.BytesIO()
        try:
            img.save(buf, format=fmt)
        except Exception:
            buf = io.BytesIO()
            fmt = "PNG"
            img.save(buf, format=fmt)
        data = buf.getvalue()

        if len(data) > MAX_BASE64_IMAGE_BYTES:
            print(f"    [base64 저장 생략] {img_path.name}: 축소 후에도 {len(data) / 1024 / 1024:.1f}MB (제한 {MAX_BASE64_IMAGE_BYTES / 1024 / 1024:.0f}MB 초과)")
            return ""
        return base64.b64encode(data).decode("utf-8")
    except Exception as e:
        print(f"    [base64 인코딩 실패] {img_path.name}: {e}")
        return ""


OCR_MAX_SIDE = 2000  # OCR은 300dpi 안팎이면 충분해서, 이보다 큰 이미지는 축소 후 OCR (정확도 손실은 미미, 속도는 크게 향상)


def ocr_image(img_path: Path) -> str:
    if not OCR_AVAILABLE:
        return ""
    try:
        img = Image.open(img_path)
        if max(img.size) > OCR_MAX_SIDE:
            scale = OCR_MAX_SIDE / max(img.size)
            new_size = (round(img.width * scale), round(img.height * scale))
            img = img.resize(new_size, Image.LANCZOS)
        return pytesseract.image_to_string(img, lang="kor+eng").strip()
    except Exception as e:
        print(f"    [OCR 실패] {img_path.name}: {e}")
        return ""


def caption_image(img_path: Path) -> str:
    """로컬 Ollama 비전 모델로 이미지의 의미를 설명하는 캡션을 생성."""
    try:
        img_b64 = base64.b64encode(img_path.read_bytes()).decode("utf-8")
        resp = requests.post(
            OLLAMA_URL,
            json={
                "model": VISION_MODEL,
                "prompt": CAPTION_PROMPT,
                "images": [img_b64],
                "stream": False,
            },
            timeout=120,
        )
        resp.raise_for_status()
        return resp.json().get("response", "").strip()
    except requests.exceptions.ConnectionError:
        # Ollama가 꺼져 있음 - 조용히 건너뜀 (파일마다 반복 출력 방지)
        return ""
    except Exception as e:
        print(f"    [캡션 생성 실패] {img_path.name}: {e}")
        return ""


MAX_IMAGES_PER_PAGE = 3  # 페이지/슬라이드/시트당 처리(추출+OCR+캡션)할 최대 이미지 수, 큰 순서대로
MAX_IMAGES_NO_PAGE = 10  # hwpx처럼 고정 페이지 개념이 없는 문서에서 처리할 전체 이미지 수 상한
MIN_IMAGE_SIDE = 80  # 이보다 작은 이미지(아이콘/장식용 라인 등)는 노이즈로 보고 건너뜀


def _report(progress_callback, unit_index: int, unit_total: int, unit_label: str):
    """파일 내부 진행률 콜백을 안전하게 호출. unit_label 예: "페이지"/"슬라이드"/"시트"/"이미지".
    콜백이 없거나 실패해도 등록 작업엔 영향 없음."""
    if progress_callback is None:
        return
    try:
        progress_callback(unit_index, unit_total, unit_label)
    except Exception:
        pass


def save_and_analyze_image(image_bytes: bytes, ext: str, img_out_dir: Path, base_name: str,
                            page: int | None, index: int) -> dict:
    """이미지 바이트를 파일로 저장하고 OCR+캡션까지 처리해서 register_images용 meta dict 생성."""
    img_out_dir.mkdir(parents=True, exist_ok=True)
    img_filename = f"{base_name}_p{page}_{index}.{ext}" if page is not None else f"{base_name}_{index}.{ext}"
    img_path = img_out_dir / img_filename
    img_path.write_bytes(image_bytes)

    ocr_text = ocr_image(img_path)
    caption = caption_image(img_path)
    status = [f"OCR {len(ocr_text)}자" if ocr_text else "OCR 없음", "캡션 O" if caption else "캡션 없음"]
    print(f"    이미지 추출: {img_filename} (" + ", ".join(status) + ")")

    return {
        "page": page, "image_index": index, "image_path": str(img_path),
        "ocr_text": ocr_text, "caption": caption,
    }


def process_pdf(path: Path, progress_callback=None) -> tuple[str, list[dict]]:
    """PDF에서 (본문 텍스트, 이미지 메타데이터 목록)을 함께 추출."""
    doc = fitz.open(str(path))
    total_pages = doc.page_count
    full_text = []
    images_meta = []
    img_out_dir = IMAGES_DIR / path.stem

    for page_num, page in enumerate(doc, start=1):
        full_text.append(page.get_text())

        if not PROCESS_IMAGES:
            _report(progress_callback, page_num, total_pages, "페이지")
            continue

        # get_images(full=True) 튜플: (xref, smask, width, height, ...)
        # 너무 작은 이미지는 제외하고, 면적(width*height) 기준 큰 순서로
        # 페이지당 MAX_IMAGES_PER_PAGE장만 골라 처리 (캡션 생성이 느려서 제한)
        candidates = [img for img in page.get_images(full=True) if img[2] >= MIN_IMAGE_SIDE and img[3] >= MIN_IMAGE_SIDE]
        candidates.sort(key=lambda img: img[2] * img[3], reverse=True)
        top_images = candidates[:MAX_IMAGES_PER_PAGE]

        for img_idx, img in enumerate(top_images, start=1):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception as e:
                print(f"    [이미지 추출 실패] p{page_num}-{img_idx}: {e}")
                continue
            meta = save_and_analyze_image(
                base_image["image"], base_image.get("ext", "png"), img_out_dir, path.stem, page_num, img_idx,
            )
            images_meta.append(meta)
            # 페이지 하나에 느린 이미지(캡션)가 여러 장 있어도 이미지 단위로 실시간 갱신
            _report(progress_callback, img_idx, len(top_images), "이미지")

        _report(progress_callback, page_num, total_pages, "페이지")

    doc.close()
    return "\n".join(full_text), images_meta


def iter_picture_shapes(shapes):
    """pptx 도형 목록을 순회하며 그림 도형만 뽑아낸다 (그룹 안에 중첩된 그림도 재귀적으로 포함)."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_picture_shapes(shape.shapes)


def process_pptx(path: Path, progress_callback=None) -> tuple[str, list[dict]]:
    """pptx에서 (본문 텍스트, 이미지 메타데이터 목록)을 함께 추출."""
    prs = Presentation(str(path))
    slides = list(prs.slides)
    total_slides = len(slides) or 1
    slides_text = []
    images_meta = []
    img_out_dir = IMAGES_DIR / path.stem

    for slide_num, slide in enumerate(slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
        if parts:
            slides_text.append(f"[슬라이드 {slide_num}]\n" + "\n".join(parts))

        if not PROCESS_IMAGES:
            _report(progress_callback, slide_num, total_slides, "슬라이드")
            continue

        candidates = []
        for shape in iter_picture_shapes(slide.shapes):
            try:
                image = shape.image
                w, h = image.size
            except Exception as e:
                print(f"    [이미지 읽기 실패] 슬라이드{slide_num}: {e}")
                continue
            if w < MIN_IMAGE_SIDE or h < MIN_IMAGE_SIDE:
                continue
            candidates.append((w * h, image.blob, image.ext))
        candidates.sort(key=lambda c: c[0], reverse=True)
        top_images = candidates[:MAX_IMAGES_PER_PAGE]

        for idx, (_, blob, ext) in enumerate(top_images, start=1):
            meta = save_and_analyze_image(blob, ext, img_out_dir, path.stem, slide_num, idx)
            images_meta.append(meta)
            _report(progress_callback, idx, len(top_images), "이미지")

        _report(progress_callback, slide_num, total_slides, "슬라이드")

    return "\n".join(slides_text), images_meta


def process_xlsx(path: Path, progress_callback=None) -> tuple[str, list[dict]]:
    """xlsx에서 (본문 텍스트, 이미지 메타데이터 목록)을 함께 추출.
    이미지는 read_only 워크북에서 접근할 수 없어 이미지 처리가 켜져 있을 때만
    일반 모드로 다시 연다."""
    img_out_dir = IMAGES_DIR / path.stem
    images_meta = []

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        lines = []
        for ws in wb.worksheets:
            lines.append(f"[시트: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        text = "\n".join(lines)
    finally:
        wb.close()

    if not PROCESS_IMAGES:
        _report(progress_callback, 1, 1, "시트")
        return text, images_meta

    wb_full = openpyxl.load_workbook(str(path), data_only=True)
    try:
        total_sheets = len(wb_full.worksheets) or 1
        for sheet_idx, ws in enumerate(wb_full.worksheets, start=1):
            candidates = []
            for img in getattr(ws, "_images", []):
                try:
                    data = img._data()
                    w, h = Image.open(io.BytesIO(data)).size
                except Exception as e:
                    print(f"    [이미지 읽기 실패] {ws.title}: {e}")
                    continue
                if w < MIN_IMAGE_SIDE or h < MIN_IMAGE_SIDE:
                    continue
                candidates.append((w * h, data))
            candidates.sort(key=lambda c: c[0], reverse=True)
            top_images = candidates[:MAX_IMAGES_PER_PAGE]

            for idx, (_, data) in enumerate(top_images, start=1):
                meta = save_and_analyze_image(data, "png", img_out_dir, path.stem, sheet_idx, idx)
                images_meta.append(meta)
                _report(progress_callback, idx, len(top_images), "이미지")

            _report(progress_callback, sheet_idx, total_sheets, "시트")
    finally:
        wb_full.close()

    return text, images_meta


def convert_with_libreoffice(path: Path, target_ext: str) -> Path | None:
    """LibreOffice로 구버전 .doc/.ppt 등을 target_ext로 변환. 실패 시 None."""
    if not SOFFICE_CMD:
        return None
    out_dir = Path(tempfile.mkdtemp(prefix="lo_convert_"))
    try:
        subprocess.run(
            [SOFFICE_CMD, "--headless", "--convert-to", target_ext, "--outdir", str(out_dir), str(path)],
            capture_output=True, timeout=120,
        )
        converted = out_dir / f"{path.stem}.{target_ext}"
        if converted.exists():
            return converted
    except Exception as e:
        print(f"    [LibreOffice 변환 실패] {path.name}: {e}")
    shutil.rmtree(out_dir, ignore_errors=True)
    return None


def read_hwpx(path: Path) -> str:
    """hwpx(zip+xml)에서 본문 텍스트만 추출."""
    texts = []
    with zipfile.ZipFile(path) as z:
        section_files = sorted(n for n in z.namelist() if re.match(r"Contents/section\d+\.xml$", n))
        for name in section_files:
            root = ET.fromstring(z.read(name))
            for elem in root.iter():
                if (elem.tag.endswith("}t") or elem.tag == "t") and elem.text:
                    texts.append(elem.text)
    return "\n".join(texts)


def process_hwpx(path: Path, progress_callback=None) -> tuple[str, list[dict]]:
    """hwpx에서 (본문 텍스트, 이미지 메타데이터 목록)을 함께 추출.
    hwpx는 워드처럼 흐르는 문서라 고정된 "페이지"가 없어, 이미지는 page=None으로
    저장하고 문서 내 순서(image_index)만 부여한다."""
    text = read_hwpx(path)
    images_meta = []
    if not PROCESS_IMAGES:
        _report(progress_callback, 1, 1, "이미지")
        return text, images_meta

    img_out_dir = IMAGES_DIR / path.stem
    candidates = []
    with zipfile.ZipFile(path) as z:
        bindata_names = sorted(n for n in z.namelist() if n.startswith("BinData/") and not n.endswith("/"))
        for name in bindata_names:
            data = z.read(name)
            try:
                w, h = Image.open(io.BytesIO(data)).size
            except Exception:
                continue  # 폰트 등 이미지가 아닌 바이너리는 건너뜀
            if w < MIN_IMAGE_SIDE or h < MIN_IMAGE_SIDE:
                continue
            ext = Path(name).suffix.lstrip(".").lower() or "png"
            candidates.append((w * h, data, ext))

    candidates.sort(key=lambda c: c[0], reverse=True)
    selected = candidates[:MAX_IMAGES_NO_PAGE]
    for idx, (_, data, ext) in enumerate(selected, start=1):
        meta = save_and_analyze_image(data, ext, img_out_dir, path.stem, None, idx)
        images_meta.append(meta)
        _report(progress_callback, idx, len(selected) or 1, "이미지")

    return text, images_meta


def read_doc(path: Path) -> str:
    try:
        out = subprocess.run(
            ["antiword", str(path)], capture_output=True, text=True, encoding="utf-8", errors="ignore"
        )
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout
    except FileNotFoundError:
        pass

    converted = convert_with_libreoffice(path, "docx")
    if converted:
        try:
            return "\n".join(p.text for p in docx.Document(str(converted)).paragraphs)
        finally:
            shutil.rmtree(converted.parent, ignore_errors=True)

    print(f"  [건너뜀] {path.name}: .doc 읽기 도구 없음 (antiword 또는 LibreOffice 설치 필요)")
    return ""


def read_hwp(path: Path) -> str:
    """pyhwp를 라이브러리로 직접 호출 (exe로 빌드 시 sys.executable이 파이썬이 아니게 되므로
    'python -m hwp5.hwp5txt' 서브프로세스 대신 이 방식을 사용)."""
    from contextlib import closing

    from hwp5.hwp5txt import TextTransform
    from hwp5.xmlmodel import Hwp5File

    transform = TextTransform().transform_hwp5_to_text
    tmp_path = Path(tempfile.mktemp(suffix=".txt"))
    try:
        with closing(Hwp5File(str(path))) as hwp5file, open(tmp_path, "wb") as dest:
            transform(hwp5file, dest)
        return tmp_path.read_text(encoding="utf-8", errors="ignore")
    finally:
        tmp_path.unlink(missing_ok=True)


def process_ppt(path: Path, progress_callback=None) -> tuple[str, list[dict]]:
    """구버전 .ppt는 LibreOffice로 .pptx 변환 후 process_pptx로 텍스트+이미지 함께 추출."""
    converted = convert_with_libreoffice(path, "pptx")
    if not converted:
        print(f"  [건너뜀] {path.name}: .ppt 변환 불가 (LibreOffice 설치 필요)")
        _report(progress_callback, 1, 1, "슬라이드")
        return "", []
    try:
        return process_pptx(converted, progress_callback=progress_callback)
    finally:
        shutil.rmtree(converted.parent, ignore_errors=True)


_HTML_BLOCK_TAGS = {
    "p", "div", "br", "li", "tr", "table", "section", "article",
    "header", "footer", "blockquote", "pre",
    "h1", "h2", "h3", "h4", "h5", "h6",
}


def read_html(path: Path) -> str:
    """HTML에서 태그/스크립트/스타일을 제거하고 실제로 보이는 본문 텍스트만 추출."""
    import lxml.html as lh

    # 바이트를 그대로 넘기면 lxml이 <meta charset>이 없는 파일에서 인코딩을 잘못
    # 추정해 한글이 깨진다. 미리 utf-8로 디코딩한 문자열을 넘겨 그 문제를 피한다.
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    tree = lh.fromstring(raw_text)
    for bad in tree.xpath("//script | //style"):
        bad.getparent().remove(bad)
    # text_content()는 블록 요소 경계에 줄바꿈을 넣어주지 않아 문단이 다 붙어버리므로,
    # 블록 태그 뒤에 직접 줄바꿈을 끼워 넣는다.
    for el in tree.iter():
        if el.tag in _HTML_BLOCK_TAGS:
            el.tail = (el.tail or "") + "\n"
    text = tree.text_content()
    # 여러 줄 공백을 하나로 정리해서 읽기 좋게
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def read_text_only(path: Path) -> str:
    """이미지 처리가 필요 없는 포맷(또는 텍스트만 필요한 호출)을 위한 텍스트 전용 추출.
    .pdf/.pptx/.ppt/.xlsx/.hwpx는 이미지까지 함께 뽑는 process_*() 함수를 register_file()에서 직접 사용."""
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            return "\n".join(p.text for p in docx.Document(str(path)).paragraphs)
        if ext == ".doc":
            return read_doc(path)
        if ext == ".hwp":
            return read_hwp(path)
        if ext in (".html", ".htm"):
            return read_html(path)
        if ext in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"  [읽기 실패] {path.name}: {e}")
    return ""


# 텍스트/이미지를 하나씩 개별 호출로 저장하면 호출마다 서버에서 임베딩을 새로 계산해야 해서
# 느리다. 이 서버의 병목은 로컬 임베딩 모델(CPU 연산)이라 동시 요청을 늘려봐야 같은 CPU를
# 두고 경쟁만 할 뿐 별 도움이 안 되고, 오히려 여러 개를 한 번에 묶어 embed_documents()를
# 한 번만 호출하는 쪽(배치)이 훨씬 효과적이다. 그래서 BATCH_SIZE개씩 묶어
# qdrant_store_batch/qdrant_store_mine_batch(서버에서 임베딩+upsert를 한 번에 처리)로 보낸다.
# 이미지는 base64까지 실리면 항목당 최대 몇 MB가 될 수 있어 요청 본문이 과도하게 커지지
# 않도록 이미지 전용으로 훨씬 작은 배치 크기를 쓴다.
TEXT_BATCH_SIZE = 25
IMAGE_BATCH_SIZE = 3
UPLOAD_CONCURRENCY = 8  # 배치 툴이 없는 게이트웨이용 폴백: 개별 저장 호출을 이만큼 동시에 전송


async def _store_group_individually(
    session: ClientSession, group: list[tuple[str, dict]], store_tool: str,
    progress_callback, unit_label: str, done: int, total: int,
) -> int:
    """배치 저장 툴이 없는(구버전) 게이트웨이용 폴백. 개별 qdrant-store 호출을
    UPLOAD_CONCURRENCY개씩 동시에 보내 순차 호출보다 훨씬 빠르게 처리한다."""
    sem = asyncio.Semaphore(UPLOAD_CONCURRENCY)
    lock = asyncio.Lock()

    async def store_one(info: str, metadata: dict):
        nonlocal done
        async with sem:
            result = await session.call_tool(store_tool, {"information": info, "metadata": metadata})
        if getattr(result, "isError", False):
            message = "\n".join(b.text for b in result.content if hasattr(b, "text")) or str(result)
            print(f"[오류] {message}")
        async with lock:
            done += 1
            _report(progress_callback, done, total, unit_label)

    await asyncio.gather(*(store_one(info, metadata) for info, metadata in group))
    return done


async def _store_in_batches(
    session: ClientSession, records: list[tuple[str, dict]], store_tool: str, batch_size: int,
    progress_callback=None, unit_label: str = "저장",
) -> list[tuple[str, dict]]:
    """(information, metadata) 목록을 batch_size개씩 묶어 배치 저장 툴로 전송.
    게이트웨이가 아직 배치 툴을 지원하지 않으면("Unknown tool") 개별 동시 저장으로 자동 전환."""
    if not records:
        return []
    batch_tool = "qdrant_store_mine_batch" if store_tool == "qdrant_store_mine" else "qdrant_store_batch"
    total = len(records)
    use_batch = True
    done = 0

    for start in range(0, total, batch_size):
        group = records[start:start + batch_size]

        if use_batch:
            items = [{"information": info, "metadata": metadata} for info, metadata in group]
            result = await session.call_tool(batch_tool, {"items": items})
            if getattr(result, "isError", False):
                message = "\n".join(b.text for b in result.content if hasattr(b, "text")) or str(result)
                if "Unknown tool" in message:
                    print(f"[안내] {batch_tool} 툴이 게이트웨이에 아직 없어 개별 동시 저장으로 전환합니다.")
                    use_batch = False
                else:
                    print(f"[오류] {message}")
                    done = min(start + batch_size, total)
                    _report(progress_callback, done, total, unit_label)
                    continue
            else:
                done = min(start + batch_size, total)
                _report(progress_callback, done, total, unit_label)
                continue

        done = await _store_group_individually(session, group, store_tool, progress_callback, unit_label, done, total)

    return records


async def register_text_chunks_raw(
    session: ClientSession, source: str, title: str, text: str,
    store_tool: str = "qdrant-store", progress_callback=None,
) -> list[tuple[str, dict]]:
    """source/title을 직접 지정해서 텍스트를 등록 (파일이 아니라 붙여넣은 텍스트 등에도 재사용).
    store_tool: "qdrant-store"(팀 공유) 또는 "qdrant_store_mine"(개인 저장소)."""
    if not text.strip():
        return []
    chunks = chunk_text(text)
    records = [
        (chunk, {"source": source, "title": title, "type": "text", "chunk_index": idx})
        for idx, chunk in enumerate(chunks)
    ]
    return await _store_in_batches(session, records, store_tool, TEXT_BATCH_SIZE, progress_callback, "저장")


async def register_text_chunks(
    session: ClientSession, path: Path, text: str, store_tool: str = "qdrant-store", progress_callback=None,
) -> list[tuple[str, dict]]:
    """Qdrant에 실제로 보낸 (information, metadata) 목록을 그대로 반환 (덤프 파일 작성에 재사용)."""
    return await register_text_chunks_raw(session, str(path), path.name, text, store_tool, progress_callback)


async def register_images(
    session: ClientSession, path: Path, images_meta: list[dict],
    store_tool: str = "qdrant-store", progress_callback=None,
) -> list[tuple[str, dict]]:
    """Qdrant에 실제로 보낸 (information, metadata) 목록을 그대로 반환 (덤프 파일 작성에 재사용)."""
    records = []
    for meta in images_meta:
        if meta["page"] is not None:
            parts = [f"[{path.name} {meta['page']}페이지에 포함된 이미지]"]
        else:
            parts = [f"[{path.name} 이미지]"]
        if meta["caption"]:
            parts.append(f"설명: {meta['caption']}")
        if meta["ocr_text"]:
            parts.append(f"이미지 내 텍스트(OCR): {meta['ocr_text']}")
        if not meta["caption"] and not meta["ocr_text"]:
            parts.append("(캡션/OCR 텍스트 없음, 이미지 파일로만 저장됨)")
        info = "\n".join(parts)

        metadata = {
            "source": str(path), "title": path.name,
            "type": "image", "page": meta["page"],
            "image_index": meta["image_index"],
            "image_path": meta["image_path"],
            "caption": meta["caption"],
        }
        if STORE_IMAGE_BASE64:
            img_b64 = encode_image_base64(Path(meta["image_path"]))
            if img_b64:
                metadata["image_base64"] = img_b64

        records.append((info, metadata))

    return await _store_in_batches(session, records, store_tool, IMAGE_BATCH_SIZE, progress_callback, "이미지 저장")


def write_extraction_dump_raw(dump_name: str, source_label: str,
                               text_records: list[tuple[str, dict]], image_records: list[tuple[str, dict]]):
    """dump_name(파일명으로 안전한 문자열)/source_label(원본 표기)을 직접 지정해서 덤프 파일 작성."""
    if not text_records and not image_records:
        return
    EXTRACTED_TEXT_DIR.mkdir(parents=True, exist_ok=True)
    lines = [f"=== 원본: {source_label} ===", ""]

    lines.append(f"--- 텍스트 청크 (총 {len(text_records)}개) ---")
    for chunk, metadata in text_records:
        lines.append(f"[청크 {metadata['chunk_index']}]")
        lines.append(chunk)
        lines.append("")

    lines.append(f"--- 이미지 (총 {len(image_records)}개) ---")
    for info, metadata in image_records:
        page = metadata["page"]
        loc = f"page={page}, index={metadata['image_index']}" if page is not None else f"index={metadata['image_index']}"
        lines.append(f"[이미지: {loc}] {metadata['image_path']}")
        lines.append(info)
        lines.append("")

    dump_path = EXTRACTED_TEXT_DIR / f"{dump_name}.txt"
    dump_path.write_text("\n".join(lines), encoding="utf-8")


def write_extraction_dump(path: Path, text_records: list[tuple[str, dict]], image_records: list[tuple[str, dict]]):
    """이 파일에서 실제로 Qdrant에 등록된 내용을 사람이 읽기 좋은 텍스트 파일로 남긴다."""
    write_extraction_dump_raw(path.stem, str(path), text_records, image_records)


# 텍스트와 이미지를 함께 뽑아내는 포맷들 (이미지가 있으면 OCR+캡션까지 처리)
TEXT_AND_IMAGE_PROCESSORS = {
    ".pdf": process_pdf,
    ".pptx": process_pptx,
    ".ppt": process_ppt,
    ".xlsx": process_xlsx,
    ".hwpx": process_hwpx,
}


async def register_file(session: ClientSession, path: Path, progress_callback=None, store_tool: str = "qdrant-store"):
    ext = path.suffix.lower()

    if ext in TEXT_AND_IMAGE_PROCESSORS:
        text, images_meta = TEXT_AND_IMAGE_PROCESSORS[ext](path, progress_callback=progress_callback)
        text_records = await register_text_chunks(session, path, text, store_tool, progress_callback)
        image_records = (
            await register_images(session, path, images_meta, store_tool, progress_callback) if images_meta else []
        )
        write_extraction_dump(path, text_records, image_records)
        print(f"등록 완료: {path.name} (텍스트 {len(text_records)}청크, 이미지 {len(image_records)}개)")
        return

    if ext in IMAGE_EXTS:
        if not PROCESS_IMAGES:
            print(f"건너뜀(이미지 처리 꺼짐): {path.name}")
            return
        ocr_text = ocr_image(path)
        caption = caption_image(path)
        meta = [{"page": None, "image_index": 1, "image_path": str(path), "ocr_text": ocr_text, "caption": caption}]
        image_records = await register_images(session, path, meta, store_tool, progress_callback)
        write_extraction_dump(path, [], image_records)
        print(f"등록 완료: {path.name} (이미지 {len(image_records)}개)")
        return

    text = read_text_only(path)
    if not text.strip():
        print(f"건너뜀(내용 없음): {path.name}")
        return
    text_records = await register_text_chunks(session, path, text, store_tool, progress_callback)
    write_extraction_dump(path, text_records, [])
    print(f"등록 완료: {path.name} (텍스트 {len(text_records)}청크)")


async def register_targets(
    session: ClientSession, targets: list[Path], progress_callback=None, store_tool: str = "qdrant-store"
):
    """파일/폴더가 섞인 경로 목록을 받아 실제 파일 목록으로 펼친 뒤 등록."""
    files = []
    for target in targets:
        if not target.exists():
            print(f"경로가 존재하지 않습니다: {target}")
            continue
        if target.is_dir():
            sub_files = [f for f in target.rglob("*") if f.is_file()]
            print(f"{target} 안 {len(sub_files)}개 파일 발견")
            files.extend(sub_files)
        else:
            files.append(target)

    if not files:
        print("등록할 파일이 없습니다.")
        return

    total = len(files)
    print(f"총 {total}개 파일 등록 시작...")
    for i, f in enumerate(files):
        def _unit_progress(unit_index, unit_total, unit_label, i=i, f=f):
            if progress_callback is None:
                return
            try:
                progress_callback({
                    "file_index": i + 1, "file_total": total, "file_name": f.name,
                    "unit_index": unit_index, "unit_total": unit_total, "unit_label": unit_label,
                })
            except Exception:
                pass

        await register_file(session, f, progress_callback=_unit_progress, store_tool=store_tool)
        _unit_progress(1, 1, "완료")  # 세부 진행률을 못 받는 포맷(.doc/.txt 등)도 파일 완료 시 확실히 갱신


async def main(targets: Path | list[Path], progress_callback=None, store_tool: str = "qdrant-store"):
    """store_tool: "qdrant-store"(팀 공유 저장소) 또는 "qdrant_store_mine"(개인 저장소)."""
    if isinstance(targets, Path):
        targets = [targets]

    if not PROCESS_IMAGES:
        print("이미지 처리: config.json의 process_images=n 설정으로 꺼져 있음 (텍스트만 등록)")
    else:
        if not OCR_LIB_AVAILABLE:
            print("참고: pytesseract/Pillow가 없어 OCR 없이 이미지만 저장합니다. (pip install pytesseract pillow)")
        elif not OCR_AVAILABLE:
            print(
                "참고: Tesseract 엔진을 찾을 수 없어 OCR 없이 진행합니다. "
                f"PATH에 없으면 config.json의 tesseract_cmd에 실제 설치 경로를 지정하세요. "
                f"(설정값: {CONFIG['tesseract_cmd']!r})"
            )

        ollama_base = OLLAMA_URL.rsplit("/api/", 1)[0]
        try:
            requests.get(f"{ollama_base}/api/tags", timeout=2)
            print(f"이미지 캡션: Ollama 연결됨 (모델: {VISION_MODEL})")
        except requests.exceptions.ConnectionError:
            print("참고: Ollama가 실행 중이 아니라 이미지 캡션 없이 진행합니다. (ollama serve 로 실행하세요)")

    if not SOFFICE_CMD:
        print("참고: LibreOffice가 없어 구버전 .ppt(및 antiword 없을 때 .doc)는 건너뜁니다.")

    async with streamablehttp_client(MCP_URL) as mcp_streams:
        # mcp 버전에 따라 (read, write) 또는 (read, write, get_session_id)를 yield하므로
        # 앞의 두 값만 안전하게 꺼내 쓴다.
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            await register_targets(session, targets, progress_callback=progress_callback, store_tool=store_tool)

    print("모든 작업 완료.")
    print(f"추출된 이미지는 여기 저장됨: {IMAGES_DIR}")
    print(f"Qdrant에 등록된 내용(텍스트+이미지 설명)은 여기서 확인 가능: {EXTRACTED_TEXT_DIR}")


def _sanitize_filename(name: str) -> str:
    name = re.sub(r'[\\/:*?"<>|]', "_", name).strip()
    return name or "pasted_text"


async def register_pasted_text(
    title: str, text: str, progress_callback=None, store_tool: str = "qdrant-store"
) -> int:
    """탐색기 파일이 아니라 GUI에 직접 붙여넣은 텍스트를 Qdrant에 등록. 등록된 청크 수 반환.
    store_tool: "qdrant-store"(팀 공유 저장소) 또는 "qdrant_store_mine"(개인 저장소)."""
    title = title.strip() or datetime.now().strftime("붙여넣은 텍스트 %Y-%m-%d %H:%M:%S")
    source = f"(직접 입력) {title}"

    async with streamablehttp_client(MCP_URL) as mcp_streams:
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            records = await register_text_chunks_raw(session, source, title, text, store_tool, progress_callback)

    write_extraction_dump_raw(_sanitize_filename(title), source, records, [])
    label = "개인 저장소" if store_tool == "qdrant_store_mine" else "팀 공유 저장소"
    print(f"등록 완료({label}): {title} (텍스트 {len(records)}청크)")
    return len(records)


async def _call_qdrant_delete(args: dict) -> int:
    """게이트웨이의 qdrant_delete 툴을 호출해 실제로 삭제된 개수를 반환하는 저수준 헬퍼."""
    async with streamablehttp_client(MCP_URL) as mcp_streams:
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            result = await session.call_tool("qdrant_delete", args)

    if getattr(result, "isError", False):
        message = "\n".join(block.text for block in result.content if hasattr(block, "text")) or str(result)
        print(f"[오류] {message}")
        return 0

    # qdrant_delete는 {"deleted": N}을 반환한다. FastMCP는 이를 structuredContent(dict)와
    # content(JSON 문자열 TextContent)로 함께 실어 보내므로, structuredContent를 우선 쓰고
    # 없으면 content의 JSON 문자열을 파싱하는 순서로 안전하게 꺼낸다.
    deleted = None
    structured = getattr(result, "structuredContent", None)
    if isinstance(structured, dict):
        deleted = structured.get("deleted")
    if deleted is None:
        for block in result.content:
            if hasattr(block, "text"):
                try:
                    deleted = json.loads(block.text).get("deleted")
                    break
                except Exception:
                    pass
    return deleted or 0


async def delete_by_source(source: str) -> int:
    """source(등록 당시 파일 경로 또는 붙여넣은 텍스트의 source 문자열)에 해당하는 모든
    항목(텍스트+이미지)을 통째로 삭제한다. 실제로 삭제된 개수를 반환.
    되돌릴 수 없는 작업이므로 호출 전 UI 쪽에서 반드시 확인을 받아야 한다."""
    deleted = await _call_qdrant_delete({"source": source})
    if deleted == 0:
        print(f"삭제할 항목이 없습니다: source={source!r}과 일치하는 데이터가 없습니다.")
    else:
        print(f"삭제 완료: source={source!r}, {deleted}개 항목을 삭제했습니다.")
    return deleted


_FIND_ENTRY_RE = re.compile(r"<entry><content>(.*?)</content><metadata>(.*?)</metadata></entry>", re.S)


async def search_qdrant(query: str) -> list[dict]:
    """qdrant-find로 검색해서 [{"content": str, "metadata": dict}, ...] 목록을 반환.
    삭제할 항목을 키워드로 찾아 고를 때 사용 (검색 자체는 아무것도 지우지 않음)."""
    async with streamablehttp_client(MCP_URL) as mcp_streams:
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            result = await session.call_tool("qdrant-find", {"query": query})

    if getattr(result, "isError", False):
        message = "\n".join(block.text for block in result.content if hasattr(block, "text")) or str(result)
        print(f"[검색 오류] {message}")
        return []

    raw = None
    structured = getattr(result, "structuredContent", None)
    if isinstance(structured, dict):
        raw = structured.get("result")
    if raw is None:
        for block in result.content:
            if hasattr(block, "text"):
                try:
                    raw = json.loads(block.text)
                    break
                except Exception:
                    pass
    if not raw:
        return []

    parsed = []
    for item in raw:
        if not isinstance(item, str):
            continue
        m = _FIND_ENTRY_RE.match(item.strip())
        if not m:
            continue  # 안내 문구("Results for the query ...") 등 entry 형식이 아닌 항목은 건너뜀
        content_text, meta_json = m.groups()
        try:
            metadata = json.loads(meta_json)
        except Exception:
            metadata = {}
        parsed.append({"content": content_text, "metadata": metadata})
    return parsed


async def delete_by_metadata(metadata: dict) -> int:
    """search_qdrant()로 얻은 metadata를 그대로 넘기면, source(+chunk_index 또는
    +page+image_index)로 정확히 그 항목 하나만 골라서 삭제한다. 파일 전체가 아니라
    검색으로 찾은 낱개 항목만 지우고 싶을 때 사용."""
    source = metadata.get("source")
    if not source:
        print("[오류] source가 없는 항목은 삭제할 수 없습니다.")
        return 0

    args = {"source": source}
    if metadata.get("type") == "image":
        if metadata.get("page") is not None:
            args["page"] = metadata["page"]
        if metadata.get("image_index") is not None:
            args["image_index"] = metadata["image_index"]
    elif metadata.get("chunk_index") is not None:
        args["chunk_index"] = metadata["chunk_index"]

    deleted = await _call_qdrant_delete(args)
    label = metadata.get("title") or source
    if deleted == 0:
        print(f"삭제할 항목이 없습니다: {label}")
    else:
        print(f"삭제 완료: {label} ({deleted}개)")
    return deleted


async def search_my_qdrant(query: str) -> list[dict]:
    """qdrant_find_mine으로 검색해서 [{"content": str, "metadata": dict}, ...] 목록을 반환.
    config.json의 mcp_url에 담긴 key 본인의 개인 저장소만 검색된다(다른 사람 데이터는
    애초에 안 보임). search_qdrant()의 개인 저장소 버전 — entry 형식이 같아서 같은
    정규식(_FIND_ENTRY_RE)으로 파싱한다. metadata에는 게이트웨이가 실어 보낸 point id가
    "_id" 키로 들어있어, delete_mine_by_metadata()가 그 항목만 정확히 지울 때 쓴다."""
    async with streamablehttp_client(MCP_URL) as mcp_streams:
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            result = await session.call_tool("qdrant_find_mine", {"query": query})

    if getattr(result, "isError", False):
        message = "\n".join(block.text for block in result.content if hasattr(block, "text")) or str(result)
        print(f"[검색 오류] {message}")
        return []

    raw = None
    structured = getattr(result, "structuredContent", None)
    if isinstance(structured, dict):
        raw = structured.get("result")
    if raw is None:
        for block in result.content:
            if hasattr(block, "text"):
                try:
                    raw = json.loads(block.text)
                    break
                except Exception:
                    pass
    if not raw:
        return []

    parsed = []
    for item in raw:
        if not isinstance(item, str):
            continue
        m = _FIND_ENTRY_RE.match(item.strip())
        if not m:
            continue  # "결과가 없습니다" 같은 안내 문구는 entry 형식이 아니라서 자동으로 걸러짐
        content_text, meta_json = m.groups()
        try:
            metadata = json.loads(meta_json)
        except Exception:
            metadata = {}
        parsed.append({"content": content_text, "metadata": metadata})
    return parsed


async def delete_mine_by_metadata(metadata: dict) -> int:
    """search_my_qdrant()로 얻은 metadata를 그대로 넘기면, 그 안의 _id(point id)로
    본인 개인 저장소에서 정확히 그 항목 하나만 삭제한다."""
    point_id = metadata.get("_id")
    if not point_id:
        print("[오류] _id가 없는 항목은 삭제할 수 없습니다.")
        return 0

    async with streamablehttp_client(MCP_URL) as mcp_streams:
        read, write = mcp_streams[0], mcp_streams[1]
        async with ClientSession(read, write) as session:
            await session.initialize()
            result = await session.call_tool("qdrant_delete_mine", {"point_ids": [point_id]})

    if getattr(result, "isError", False):
        message = "\n".join(block.text for block in result.content if hasattr(block, "text")) or str(result)
        print(f"[오류] {message}")
        return 0

    deleted = None
    structured = getattr(result, "structuredContent", None)
    if isinstance(structured, dict):
        deleted = structured.get("deleted")
    if deleted is None:
        for block in result.content:
            if hasattr(block, "text"):
                try:
                    deleted = json.loads(block.text).get("deleted")
                    break
                except Exception:
                    pass
    deleted = deleted or 0

    label = metadata.get("title") or metadata.get("source") or point_id
    if deleted == 0:
        print(f"삭제할 항목이 없습니다: {label}")
    else:
        print(f"삭제 완료(개인 저장소): {label}")
    return deleted


if __name__ == "__main__":
    if len(sys.argv) > 1:
        target_path = Path(sys.argv[1]).resolve()
    else:
        target_path = DEFAULT_INCOMING_DIR
        target_path.mkdir(exist_ok=True)
        print(f"인자 없이 실행됨 -> 기본 폴더 사용: {target_path}")

    asyncio.run(main(target_path))